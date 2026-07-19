import asyncio
import io
import logging
import datetime
from typing import Annotated, List
import uuid
from fastapi import APIRouter, BackgroundTasks, HTTPException, UploadFile, File
from pydantic import BaseModel

from core.rag_manager import RAGManager
from core.mongo_manager import mongo_db

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/companies", tags=["companies"])
rag_manager = RAGManager()

# In-memory job progress store: {job_id: {status, progress, total, message, filename}}
_job_store: dict = {}

def extract_text_from_file(filename: str, body: bytes) -> str:
    """Helper to extract text from raw bytes based on extension (PDF/DOCX/PPTX/Text)"""
    fn = filename.lower()
    if fn.endswith('.pdf'):
        try:
            import PyPDF2
            pdf_file = io.BytesIO(body)
            reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page in reader.pages:
                text += (page.extract_text() or "") + "\n"
            return text
        except Exception as e:
            logger.error(f"Failed to parse PDF {filename}: {e}")
            raise ValueError(f"Could not parse PDF document: {str(e)}")
    elif fn.endswith('.docx'):
        try:
            import docx
            doc = docx.Document(io.BytesIO(body))
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            return text
        except Exception as e:
            logger.error(f"Failed to parse DOCX {filename}: {e}")
            raise ValueError(f"Could not parse DOCX document: {str(e)}")
    elif fn.endswith('.pptx') or fn.endswith('.ppt'):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(body))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Failed to parse PPTX {filename}: {e}")
            raise ValueError(f"Could not parse PPTX presentation: {str(e)}")
    else:
        try:
            return body.decode('utf-8', errors='ignore')
        except Exception as e:
            logger.error(f"Failed to decode text file {filename}: {e}")
            raise ValueError(f"Could not decode file as text: {str(e)}")


@router.post("/", summary="Create a new company")
async def create_company(name: str, phone_number: str):
    if not mongo_db.client:
        raise HTTPException(status_code=503, detail="MongoDB connection offline")

    # Clean phone number
    cleaned_phone = phone_number.replace("+", "").strip()
    
    db = mongo_db.client.get_default_database()
    companies_collection = db['companies']
    
    # Check if phone number already exists
    existing = await companies_collection.find_one({"phone_number": cleaned_phone})
    if existing:
        raise HTTPException(status_code=400, detail=f"Phone number {phone_number} is already registered to company {existing.get('name')}")
        
    company_id = str(uuid.uuid4())
    company = {
        "company_id": company_id,
        "name": name,
        "phone_number": cleaned_phone,
        "metadata_json": {},
        "created_at": datetime.datetime.utcnow().isoformat() + "Z"
    }
    
    try:
        await companies_collection.insert_one(company)
        # Initialize tenant Chroma collection
        rag_manager.init_index(company_id)
    except Exception as e:
        logger.error(f"Failed to create company: {e}")
        raise HTTPException(status_code=500, detail=f"Database error: {str(e)}")
        
    return {"company_id": company_id, "name": name, "phone_number": cleaned_phone}


@router.get("/", summary="List all companies")
async def list_companies():
    if not mongo_db.client:
        raise HTTPException(status_code=503, detail="MongoDB connection offline")

    db = mongo_db.client.get_default_database()
    companies_collection = db['companies']
    
    cursor = companies_collection.find({})
    companies = []
    async for c in cursor:
        companies.append({
            "company_id": c.get("company_id"),
            "name": c.get("name"),
            "phone_number": c.get("phone_number")
        })
    return companies


@router.get("/{company_id}", summary="Get company details")
async def get_company(company_id: str):
    if not mongo_db.client:
        raise HTTPException(status_code=503, detail="MongoDB connection offline")

    db = mongo_db.client.get_default_database()
    companies_collection = db['companies']
    document_logs_collection = db['document_logs']
    
    company = await companies_collection.find_one({"company_id": company_id})
    if not company:
        raise HTTPException(status_code=404, detail="Company not found")
    
    # Fetch document logs
    cursor = document_logs_collection.find({"company_id": company_id})
    documents = []
    async for d in cursor:
        documents.append({
            "id": d.get("id"),
            "filename": d.get("filename"),
            "size_bytes": d.get("size_bytes"),
            "uploaded_at": d.get("uploaded_at"),
            "status": d.get("status")
        })
    
    return {
        "company_id": company.get("company_id"),
        "name": company.get("name"),
        "phone_number": company.get("phone_number"),
        "documents": documents
    }


@router.delete("/{company_id}", summary="Delete a company and its RAG data")
async def delete_company(company_id: str):
    if not mongo_db.client:
        raise HTTPException(status_code=503, detail="MongoDB connection offline")

    db = mongo_db.client.get_default_database()
    companies_collection = db['companies']
    document_logs_collection = db['document_logs']
    
    company = await companies_collection.find_one({"company_id": company_id})
    if not company:
        raise HTTPException(status_code=404, detail="Company not found")
        
    try:
        # Delete document logs
        await document_logs_collection.delete_many({"company_id": company_id})
        # Delete company registry
        await companies_collection.delete_one({"company_id": company_id})
        # Delete index from Chroma
        rag_manager.delete_index(company_id)
    except Exception as e:
        logger.error(f"Failed to delete company {company_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Database error: {str(e)}")
        
    return {"deleted": True, "company_id": company_id}


async def _process_document_bg(job_id: str, company_id: str, doc_log_id: str, filename: str, body: bytes, text: str):
    """Background task: embed chunks and index into Chroma, updating job progress."""
    if not mongo_db.client:
        logger.error("MongoDB offline inside background processing task")
        return

    db = mongo_db.client.get_default_database()
    document_logs_collection = db['document_logs']
    
    try:
        def on_progress(current: int, total: int):
            _job_store[job_id].update({
                "progress": current,
                "total": total,
                "message": f"Indexing chunk {current} of {total}..."
            })

        await rag_manager.upload_documents(company_id, filename, body, text, doc_log_id, on_progress=on_progress)

        await document_logs_collection.update_one(
            {"id": doc_log_id},
            {"$set": {"status": "processed"}}
        )

        _job_store[job_id].update({"status": "done", "message": "Document indexed successfully!"})
        logger.info(f"Background job {job_id}: completed indexing '{filename}' (ID: {doc_log_id})")

    except Exception as e:
        logger.error(f"Background job {job_id}: failed indexing '{filename}': {e}")
        _job_store[job_id].update({"status": "error", "message": str(e)})
        await document_logs_collection.update_one(
            {"id": doc_log_id},
            {"$set": {"status": "failed"}}
        )


@router.post(
    "/{company_id}/documents",
    summary="Upload documents for a company",
    description="Upload one or more documents (PDF, DOCX, PPTX, TXT) to index into the agent's knowledge base. Returns document IDs that can be linked to an agent via `knowledgeBaseIds`.",
    openapi_extra={
        "requestBody": {
            "content": {
                "multipart/form-data": {
                    "schema": {
                        "type": "object",
                        "properties": {
                            "files": {
                                "type": "array",
                                "items": {
                                    "type": "string",
                                    "format": "binary"
                                },
                                "description": "Select one or more files to upload (PDF, DOCX, PPTX, TXT)"
                            }
                        },
                        "required": ["files"]
                    }
                }
            },
            "required": True
        }
    }
)
async def upload_documents(
    company_id: str,
    background_tasks: BackgroundTasks,
    files: Annotated[List[UploadFile], File(description="Select one or more files (PDF, DOCX, PPTX, TXT) to upload and index.")]
):
    if not mongo_db.client:
        raise HTTPException(status_code=503, detail="MongoDB connection offline")

    db = mongo_db.client.get_default_database()
    companies_collection = db['companies']
    document_logs_collection = db['document_logs']
    
    company = await companies_collection.find_one({"company_id": company_id})
    if not company:
        raise HTTPException(status_code=404, detail="Company not found")

    started_jobs = []

    for upload in files:
        # Check if file with same name has already been uploaded for this company
        existing_doc = await document_logs_collection.find_one({
            "company_id": company_id,
            "filename": upload.filename,
            "status": {"$ne": "failed"}
        })
        if existing_doc:
            raise HTTPException(
                status_code=400,
                detail=f"Document '{upload.filename}' already exists or is being processed. Please delete it first if you want to re-upload it."
            )

        body = await upload.read()
        try:
            text = extract_text_from_file(upload.filename, body)
        except ValueError as val_err:
            raise HTTPException(status_code=400, detail=str(val_err))

        doc_id = str(uuid.uuid4())
        doc_log = {
            "id": doc_id,
            "company_id": company_id,
            "filename": upload.filename,
            "size_bytes": len(body),
            "status": "processing",
            "uploaded_at": datetime.datetime.utcnow().isoformat() + "Z"
        }
        await document_logs_collection.insert_one(doc_log)

        # Register job in progress store
        job_id = str(uuid.uuid4())
        _job_store[job_id] = {
            "status": "processing",
            "filename": upload.filename,
            "progress": 0,
            "total": 0,
            "message": "Uploading to S3 and extracting text..."
        }

        # Schedule background embedding (returns to browser immediately)
        background_tasks.add_task(
            _process_document_bg,
            job_id, company_id, doc_id, upload.filename, body, text
        )

        started_jobs.append({
            "job_id": job_id,
            "doc_id": doc_id,
            "filename": upload.filename
        })

    return {"status": "processing", "company_id": company_id, "jobs": started_jobs}


@router.get("/jobs/{job_id}", summary="Poll document indexing job status")
async def get_job_status(job_id: str):
    """Poll this endpoint to get live embedding progress for a document upload job."""
    job = _job_store.get(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job not found")
    return job


@router.get("/{company_id}/search", summary="Search the RAG store for a company")
async def search_company(company_id: str, q: str, top_k: int = 3):
    if not mongo_db.client:
        raise HTTPException(status_code=503, detail="MongoDB connection offline")

    db = mongo_db.client.get_default_database()
    companies_collection = db['companies']
    
    company = await companies_collection.find_one({"company_id": company_id})
    if not company:
        raise HTTPException(status_code=404, detail="Company not found")
        
    results = await rag_manager.search(company_id, q, top_k)
    return [{"chunk": r["chunk_text"], "source": r["metadata"].get("source")} for r in results]


@router.delete("/{company_id}/documents/{document_id}", summary="Delete a document and its RAG vectors")
async def delete_document(company_id: str, document_id: str):
    if not mongo_db.client:
        raise HTTPException(status_code=503, detail="MongoDB connection offline")

    db = mongo_db.client.get_default_database()
    companies_collection = db['companies']
    document_logs_collection = db['document_logs']
    
    company = await companies_collection.find_one({"company_id": company_id})
    if not company:
        raise HTTPException(status_code=404, detail="Company not found")
        
    doc = await document_logs_collection.find_one({"id": document_id, "company_id": company_id})
    if not doc:
        raise HTTPException(status_code=404, detail="Document log not found")
        
    try:
        # Delete from S3 and Chroma using document_id to isolate deletes
        rag_manager.delete_document(company_id, document_id, doc.get("filename"))
        
        # Delete database log
        await document_logs_collection.delete_one({"id": document_id})
    except Exception as e:
        logger.error(f"Failed to delete document {document_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Database error: {str(e)}")
        
    return {"deleted": True, "document_id": document_id, "filename": doc.get("filename")}


class RawTextInput(BaseModel):
    text: str
    source_name: str


@router.post("/{company_id}/webpages", summary="Crawl and index a web page")
async def index_webpage(company_id: str, url: str):
    if not mongo_db.client:
        raise HTTPException(status_code=503, detail="MongoDB connection offline")

    db = mongo_db.client.get_default_database()
    companies_collection = db['companies']
    document_logs_collection = db['document_logs']
    
    company = await companies_collection.find_one({"company_id": company_id})
    if not company:
        raise HTTPException(status_code=404, detail="Company not found")
        
    try:
        from playwright.async_api import async_playwright
        from bs4 import BeautifulSoup
        
        logger.info(f"Crawl URL: {url} using Playwright headless Chromium...")
        async with async_playwright() as p:
            browser = await p.chromium.launch(
                headless=True,
                args=["--no-sandbox", "--disable-setuid-sandbox"]
            )
            context = await browser.new_context(
                user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
            )
            page = await context.new_page()
            
            try:
                await page.goto(url, wait_until="networkidle", timeout=15000)
            except Exception as navigation_err:
                logger.warning(f"Playwright navigation warning for {url}: {navigation_err}")
                await page.wait_for_load_state("domcontentloaded")
                
            html = await page.content()
            await browser.close()
            
        soup = BeautifulSoup(html, 'html.parser')
        
        for script in soup(["script", "style", "nav", "footer", "header"]):
            script.decompose()
            
        text = soup.get_text()
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        clean_text = "\n".join(chunk for chunk in chunks if chunk)
        
        if not clean_text.strip():
            raise ValueError("Cleaned webpage content is empty.")
            
    except Exception as e:
        logger.error(f"Failed to fetch or parse webpage {url} via Playwright: {e}")
        raise HTTPException(status_code=400, detail=f"Webpage fetch failed: {str(e)}")
        
    filename = url.replace("https://", "").replace("http://", "").replace("/", "_").replace("?", "_").replace("&", "_")
    if len(filename) > 100:
        filename = filename[:97] + "..."
    filename = f"webpage_{filename}.txt"
    
    doc_id = str(uuid.uuid4())
    doc_log = {
        "id": doc_id,
        "company_id": company_id,
        "filename": filename,
        "size_bytes": len(clean_text.encode('utf-8')),
        "status": "processing",
        "uploaded_at": datetime.datetime.utcnow().isoformat() + "Z"
    }
    await document_logs_collection.insert_one(doc_log)
    
    try:
        await rag_manager.upload_documents(company_id, filename, clean_text.encode('utf-8'), clean_text, doc_id)
        await document_logs_collection.update_one(
            {"id": doc_id},
            {"$set": {"status": "processed"}}
        )
    except Exception as e:
        logger.error(f"Failed to index webpage {url}: {e}")
        await document_logs_collection.update_one(
            {"id": doc_id},
            {"$set": {"status": "failed"}}
        )
        raise HTTPException(status_code=500, detail=f"Failed to index webpage: {str(e)}")
        
    return {"status": "success", "company_id": company_id, "filename": filename, "source": url}


@router.post("/{company_id}/text", summary="Index raw text messages or custom instructions")
async def index_raw_text(company_id: str, payload: RawTextInput):
    if not mongo_db.client:
        raise HTTPException(status_code=503, detail="MongoDB connection offline")

    db = mongo_db.client.get_default_database()
    companies_collection = db['companies']
    document_logs_collection = db['document_logs']
    
    company = await companies_collection.find_one({"company_id": company_id})
    if not company:
        raise HTTPException(status_code=404, detail="Company not found")
        
    text_content = payload.text.strip()
    if not text_content:
        raise HTTPException(status_code=400, detail="Text content cannot be empty")
        
    source_name = payload.source_name.strip()
    if not source_name.endswith('.txt'):
        source_name = f"{source_name}.txt"
        
    import time
    timestamp = int(time.time())
    filename = f"text_{timestamp}_{source_name}"
    
    doc_id = str(uuid.uuid4())
    doc_log = {
        "id": doc_id,
        "company_id": company_id,
        "filename": filename,
        "size_bytes": len(text_content.encode('utf-8')),
        "status": "processing",
        "uploaded_at": datetime.datetime.utcnow().isoformat() + "Z"
    }
    await document_logs_collection.insert_one(doc_log)
    
    try:
        await rag_manager.upload_documents(company_id, filename, text_content.encode('utf-8'), text_content, doc_id)
        await document_logs_collection.update_one(
            {"id": doc_id},
            {"$set": {"status": "processed"}}
        )
    except Exception as e:
        logger.error(f"Failed to index raw text {filename}: {e}")
        await document_logs_collection.update_one(
            {"id": doc_id},
            {"$set": {"status": "failed"}}
        )
        raise HTTPException(status_code=500, detail=f"Failed to index raw text: {str(e)}")
        
    return {"status": "success", "company_id": company_id, "filename": filename}
